"""Build the INIT'26 submission deck from the official template.

    ./.venv/Scripts/python.exe scripts/build_deck.py

Re-runnable on purpose: the screenshots in ``website images/`` get re-shot
whenever the deployed app changes, and the deck must be rebuildable from the
current ones rather than hand-edited each time.

**Template rules, enforced here rather than trusted:**
  - exactly 8 slides, no additions, no removals
  - the seven section headings are never touched
  - nothing is placed outside the template's own content area

The script asserts the slide count before and after, so a change that would
break the submission rules fails loudly instead of producing an invalid deck.

Every number written onto a slide comes from a real run — see
`scripts/demo_figures.py`, which regenerates them, and `docs/10-RULES.md`
section 5.3, which forbids putting an unverified figure on a slide.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "INIT'26 PPT FORMAT.pptx"
OUTPUT = ROOT / "CCE — INIT26 Submission.pptx"
SHOTS = ROOT / "website images"

# ---------------------------------------------------------------------------
# Design tokens, sampled from the template rather than invented.
# ---------------------------------------------------------------------------

GREEN = RGBColor(0x2D, 0xAF, 0x4B)     # the template's own chevron accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB4, 0xBD, 0xC6)
DIM = RGBColor(0x84, 0x8D, 0x99)
AMBER = RGBColor(0xF5, 0xA5, 0x24)
RED = RGBColor(0xE5, 0x5B, 0x4C)

#: Headings are Antonio Bold. Body copy is not: Antonio is a condensed display
#: face, and a paragraph set in it is markedly harder to read at distance.
#: Calibri ships with Office on every machine that will open this file.
BODY = "Calibri"

#: The template's own content column: headings sit at x=1.61 and the slide is
#: 20.62in wide. Staying inside this is what "do not modify the format" means
#: in practice.
LEFT = Inches(1.61)
TOP = Inches(2.15)
WIDTH = Inches(17.40)
BOTTOM = Inches(10.70)

EXPECTED_SLIDES = 8


def main() -> int:
    if not TEMPLATE.exists():
        print(f"template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    prs = Presentation(str(TEMPLATE))
    if len(prs.slides) != EXPECTED_SLIDES:
        print(
            f"template has {len(prs.slides)} slides, expected "
            f"{EXPECTED_SLIDES} — refusing to guess at the structure",
            file=sys.stderr,
        )
        return 1

    _title(prs.slides[1])
    _solution(prs.slides[2])
    _technical(prs.slides[3])
    _flow(prs.slides[4])
    missing = _snapshots(prs.slides[5])
    _impact(prs.slides[6])
    _references(prs.slides[7])

    assert len(prs.slides) == EXPECTED_SLIDES, "slide count changed — rule broken"
    prs.save(str(OUTPUT))

    print(f"\nwrote {OUTPUT.name}")
    print(f"  {len(prs.slides)} slides (template limit is {EXPECTED_SLIDES})")
    if missing:
        print("\n  MISSING SCREENSHOTS — placeholders drawn instead:")
        for name in missing:
            print(f"    {name}")
    return 0


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _box(slide, left, top, width, height):
    """A plain text box with no autofit surprises."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text="", *, size=22, bold=False, colour=WHITE, space_before=0,
       space_after=6, first=False, align=PP_ALIGN.LEFT, italic=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    if text:
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = colour
        run.font.name = BODY
    return para


def _rich(tf, parts, *, size=22, space_before=0, space_after=6, first=False):
    """One paragraph built from (text, colour, bold) tuples."""
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)
    for text, colour, bold in parts:
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = BODY
    return para


def _label(tf, text, *, first=False):
    """A small green section label."""
    return _p(tf, text.upper(), size=15, bold=True, colour=GREEN,
              space_before=0 if first else 16, space_after=6, first=first)


# ---------------------------------------------------------------------------
# SLIDE 1 — title
# ---------------------------------------------------------------------------


def _title(slide) -> None:
    """Fill the four template fields.

    The template sets these at 85pt, which fits four bare labels and nothing
    else. Values are appended at a size that actually fits the box; the font,
    colour and the labels themselves are unchanged.
    """
    values = {
        "TRACK": "FinTech — Asset & Capital Management / Optimization Controls",
        "PROBLEM STATEMENT": (
            "Automate capital allocation while independently enforcing risk "
            "control"
        ),
        "TEAM NAME": "Qoders",
        "IDEA TITLE": "CCE — Capital Control Engine",
    }

    target = None
    for shape in slide.shapes:
        if shape.has_text_frame and "TRACK" in shape.text_frame.text:
            target = shape
            break
    if target is None:
        return

    tf = target.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, (label, value) in enumerate(values.items()):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(20)
        head = para.add_run()
        head.text = f"{label} : "
        head.font.size = Pt(30)
        head.font.bold = True
        head.font.color.rgb = WHITE
        head.font.name = "Antonio Bold"

        body = para.add_run()
        body.text = value
        body.font.size = Pt(30)
        body.font.bold = False
        body.font.color.rgb = GREEN if label == "IDEA TITLE" else MUTED
        body.font.name = BODY

    tag = _box(slide, LEFT, Inches(8.30), WIDTH, Inches(2.0))
    _rich(
        tag,
        [("Optimal ", WHITE, True), ("≠", GREEN, True), (" Safe.", WHITE, True)],
        size=40, first=True,
    )
    _p(
        tag,
        "The highest-return mathematical allocation is not automatically the "
        "allocation an institution should accept.",
        size=19, colour=DIM, space_before=4,
    )


# ---------------------------------------------------------------------------
# SLIDE 2 — proposed solution
# ---------------------------------------------------------------------------


def _solution(slide) -> None:
    tf = _box(slide, LEFT, TOP, WIDTH, BOTTOM - TOP)

    _rich(
        tf,
        [
            ("The optimizer proposes. An independent control engine ", WHITE, True),
            ("disposes.", GREEN, True),
        ],
        size=34, space_after=14, first=True,
    )
    _p(
        tf,
        "A ₹100 Cr institutional book on Indian market data. CCE continuously "
        "re-evaluates the portfolio, and refuses to adopt an optimizer output "
        "that breaches policy — however attractive its numbers.",
        size=21, colour=MUTED, space_after=4,
    )

    _label(tf, "Two failure modes, not one")
    _p(tf, "Most tools solve one and ignore the other.", size=18, colour=DIM,
       space_after=8)
    _rich(tf, [
        ("Stale allocation — ", WHITE, True),
        ("static rules do not adapt when market conditions change.", MUTED, False),
    ], size=20, space_after=4)
    _rich(tf, [
        ("Unsafe allocation — ", WHITE, True),
        ("a pure optimizer produces a statistically attractive portfolio that "
         "violates institutional policy.", MUTED, False),
    ], size=20)

    _label(tf, "What the control engine does that a constraint solver does not")
    for text in (
        "Re-derives every metric from raw returns. It never reads the "
        "optimizer's own report of its output, so a solver bug becomes a "
        "REJECTION rather than an approval.",
        "Judges risk contribution, not just weight. A position can sit inside "
        "its 30% weight cap while causing 62% of portfolio risk.",
        "On failure, does less — never something different. The Last Approved "
        "Safe Allocation is preserved; the system never invents one and never "
        "relaxes a limit to produce an answer.",
        "Requires a human to approve. Nothing is adopted automatically, and "
        "every approval and rejection is recorded.",
    ):
        _rich(tf, [("▸  ", GREEN, True), (text, MUTED, False)],
              size=20, space_after=7)


# ---------------------------------------------------------------------------
# SLIDE 3 — technical approach
# ---------------------------------------------------------------------------


def _technical(slide) -> None:
    half = Inches(8.45)
    gap = Inches(0.50)

    left = _box(slide, LEFT, TOP, half, BOTTOM - TOP)
    _label(left, "Layered architecture — enforced by tests, not convention",
           first=True)
    for layer, note in (
        ("ui/", "may import ONLY cce.services + cce.contracts"),
        ("cce/services/", "orchestrates; the only layer the UI touches"),
        ("cce/risk/", "pure functions, no I/O"),
        ("cce/optimizer/", "proposes only; never writes state"),
        ("cce/controls/", "INDEPENDENT authority; cannot import the optimizer"),
        ("cce/stress/", "independent gate"),
        ("cce/backtest/", "walk-forward, look-ahead prevention"),
        ("cce/audit/", "append-only SQLite; the ONLY database access"),
    ):
        _rich(left, [(f"{layer:<16}", GREEN, True), (note, MUTED, False)],
              size=17, space_after=5)

    _p(
        left,
        "tests/test_architecture.py parses every import in the tree. A layer "
        "violation fails the build rather than being caught in review — it "
        "caught two real ones during development.",
        size=17, colour=DIM, space_before=10,
    )

    right = _box(slide, LEFT + half + gap, TOP, half, BOTTOM - TOP)
    _label(right, "Financial methods", first=True)
    _p(
        right,
        "Mean-variance optimization · EWMA volatility · historical VaR and "
        "CVaR · risk contribution · concentration, liquidity and turnover "
        "constraints · transaction costs · Rockafellar–Uryasev CVaR LP · "
        "Hierarchical Risk Parity · Black-Litterman · walk-forward backtesting.",
        size=19, colour=MUTED,
    )

    _label(right, "Correctness, checked against identities")
    for text in (
        "Σ RCᵢ = σₚ to 1e-12 — a free check on the whole covariance and "
        "weight pipeline",
        "CVaR ≥ VaR across 20 seeds",
        "Annualisation applied exactly once — a √252 applied twice is a "
        "15.9× error that still looks like a number",
        "√(w'Σw) and the return series' own σ both give 10.35% — two "
        "independent paths agreeing",
    ):
        _rich(right, [("▸  ", GREEN, True), (text, MUTED, False)],
              size=17, space_after=5)

    _label(right, "Stack")
    _p(
        right,
        "Python 3.11 · NumPy · pandas · SciPy · CVXPY · jugaad-data (NSE/RBI) "
        "· SQLite · Streamlit + Plotly · pytest, ruff, mypy",
        size=17, colour=MUTED,
    )
    _rich(
        right,
        [
            ("575 tests, none skipped.", WHITE, True),
            ("  ruff and mypy clean across 93 files.  All 12 safety "
             "invariants have a real test.", MUTED, False),
        ],
        size=18, space_before=10,
    )


# ---------------------------------------------------------------------------
# SLIDE 4 — how it works
# ---------------------------------------------------------------------------


def _flow(slide) -> None:
    tf = _box(slide, LEFT, TOP, WIDTH, Inches(1.4))
    _rich(
        tf,
        [
            ("Detect → Optimize → Validate → Stress-test → Explain → ", WHITE, True),
            ("Human approval", GREEN, True),
            (" → Audit", WHITE, True),
        ],
        size=30, first=True,
    )
    _p(
        tf,
        "The optimizer is deliberately not the final authority. Each step "
        "below has an owning module, a defined failure behaviour, and no step "
        "may be performed by the UI.",
        size=18, colour=DIM, space_before=6,
    )

    steps = [
        ("1  DETECT", "cce/risk/",
         "A scheduled or user-triggered re-read of the book. Volatility, VaR, "
         "CVaR, drawdown and risk contribution are recomputed from the price "
         "panel."),
        ("2  OPTIMIZE", "cce/optimizer/",
         "Constrained max-Sharpe proposes an allocation. An unconstrained "
         "optimum is solved alongside it so the two can be shown side by side."),
        ("3  VALIDATE", "cce/controls/",
         "The control engine re-derives every metric itself and judges the "
         "candidate against 15 configured controls. It cannot import the "
         "optimizer."),
        ("4  STRESS", "cce/stress/",
         "Seven configured scenarios plus custom shocks. A scenario that could "
         "not run is reported as such — never as one the book survived."),
        ("5  BREAKER", "cce/controls/",
         "Any hard control at RED trips the circuit breaker. The Last Approved "
         "Safe Allocation is preserved and up to three independently validated "
         "recovery allocations are generated."),
        ("6  EXPLAIN", "cce/decisions/",
         "A deterministic narrator turns the structured decision record into "
         "prose. An optional LLM rephrases it — and can never alter a weight, "
         "threshold, score or approval."),
        ("7  APPROVE", "cce/services/",
         "A human approves, rejects, or keeps the current allocation. The gate "
         "is enforced server-side, not by disabling a button."),
        ("8  AUDIT", "cce/audit/",
         "Every event, candidate, finding and human action is written to an "
         "append-only store, and replayable as a timeline."),
    ]

    col_w = Inches(8.45)
    for i, (num, module, text) in enumerate(steps):
        col, row = i // 4, i % 4
        box = _box(
            slide,
            LEFT + col * (col_w + Inches(0.50)),
            Inches(4.30) + row * Inches(1.62),
            col_w, Inches(1.50),
        )
        _rich(box, [(num, GREEN, True), (f"    {module}", DIM, False)],
              size=19, space_after=3, first=True)
        _p(box, text, size=17, colour=MUTED)


# ---------------------------------------------------------------------------
# SLIDE 5 — project snapshots
# ---------------------------------------------------------------------------

#: (file, caption). Cropping is by fraction of height so a re-shot screenshot
#: of a different size still crops to roughly the same region.
PANELS = [
    ("op1.png", "Safe vs Optimal — the rejected optimum, with the specific "
                "limits it broke", None),
    ("port2.png", "Weight vs risk contribution — where the money is, and "
                  "where the risk is", (0.0, 0.03, 1.0, 0.36)),
    ("back1.png", "Walk-forward backtest — controlled vs uncontrolled on "
                  "identical data", None),
    ("decis.png", "Decision replay — every event, candidate and verdict, "
                  "reconstructed from stored records", (0.0, 0.0, 1.0, 0.46)),
]


def _find(name: str):
    """Locate a screenshot, tolerating stray spaces in the filename.

    Screenshots are hand-captured and land here with whatever name the OS
    dialog produced — "port2 .png" cost one silent placeholder before this
    existed. Matching on the squeezed stem is cheaper than asking someone to
    rename files under time pressure.
    """
    want = name.replace(" ", "").lower()
    exact = SHOTS / name
    if exact.exists():
        return exact
    for candidate in SHOTS.glob("*"):
        if candidate.name.replace(" ", "").lower() == want:
            return candidate
    return None


def _snapshots(slide) -> list[str]:
    from PIL import Image

    intro = _box(slide, LEFT, Inches(1.95), WIDTH, Inches(0.55))
    _p(
        intro,
        "Live at capitalcontrol-engine.streamlit.app — runs with no network "
        "and no API key.",
        size=19, colour=MUTED, first=True,
    )

    build = ROOT / "build" / "deck"
    build.mkdir(parents=True, exist_ok=True)

    # Sized so the SECOND row's caption still lands above the bottom edge:
    # origin + 2*(cell + gap) + caption must stay under 11.25in. The first
    # attempt overflowed by 0.16in, which a bounds check caught and an eye
    # would not have.
    cell_w, cell_h = Inches(8.45), Inches(3.25)
    gap_x, gap_y = Inches(0.50), Inches(0.78)
    origin_y = Inches(2.62)

    missing: list[str] = []
    for i, (name, caption, crop) in enumerate(PANELS):
        col, row = i % 2, i // 2
        x = LEFT + col * (cell_w + gap_x)
        y = origin_y + row * (cell_h + gap_y)

        src = _find(name)
        if src is None:
            missing.append(name)
            ph = _box(slide, x, y, cell_w, Inches(0.6))
            _p(ph, f"[ {name} — not found in 'website images/' ]",
               size=18, colour=AMBER, first=True)
        else:
            path = src
            if crop:
                img = Image.open(src)
                w, h = img.size
                x0, y0, x1, y1 = crop
                img = img.crop(
                    (int(x0 * w), int(y0 * h), int(x1 * w), int(y1 * h))
                )
                path = build / f"crop_{name}"
                img.save(path)

            img = Image.open(path)
            ratio = img.height / img.width
            draw_w = cell_w
            draw_h = Emu(int(draw_w * ratio))
            if draw_h > cell_h:                       # fit by height instead
                draw_h = cell_h
                draw_w = Emu(int(draw_h / ratio))
            slide.shapes.add_picture(
                str(path), x + Emu(int((cell_w - draw_w) / 2)), y,
                width=draw_w, height=draw_h,
            )

        cap = _box(slide, x, y + cell_h + Inches(0.04), cell_w, Inches(0.52))
        _rich(cap, [("▸  ", GREEN, True), (caption, MUTED, False)],
              size=16, first=True)

    return missing


# ---------------------------------------------------------------------------
# SLIDE 6 — feasibility, viability, impact
# ---------------------------------------------------------------------------


def _impact(slide) -> None:
    tf = _box(slide, LEFT, TOP, WIDTH, Inches(2.9))
    _label(tf, "Does the control layer actually help?", first=True)
    _p(
        tf,
        "Walk-forward, Sep 2024 – Aug 2026, monthly rebalance, identical data "
        "down both arms. Every rebalance uses only data strictly before that "
        "date; a test injects a one-day leak and fails if it goes undetected.",
        size=18, colour=DIM, space_after=10,
    )

    rows = [
        ("Strategy", "Return", "Volatility", "Max drawdown", "Policy breaches",
         DIM, True),
        ("Buy and hold", "12.8%", "8.7%", "8.8%", "0", MUTED, False),
        ("Uncontrolled optimizer", "33.5%", "11.3%", "6.5%", "37", MUTED, False),
        ("CCE-controlled", "23.5%", "7.3%", "4.8%", "0", WHITE, True),
    ]
    xs = [LEFT, LEFT + Inches(5.6), LEFT + Inches(7.7),
          LEFT + Inches(10.1), LEFT + Inches(13.0)]
    for r, (*cells, colour, bold) in enumerate(rows):
        for c, cell in enumerate(cells):
            w = Inches(5.4) if c == 0 else Inches(2.6)
            cbox = _box(slide, xs[c], Inches(3.55) + r * Inches(0.52), w,
                        Inches(0.48))
            hue = colour
            if r == 3 and c in (2, 3, 4):
                hue = GREEN
            if r == 2 and c == 4:
                hue = RED
            _p(cbox, cell, size=20, bold=bold, colour=hue, first=True)

    honest = _box(slide, LEFT, Inches(5.75), WIDTH, Inches(1.0))
    _rich(
        honest,
        [
            ("The controlled strategy earned ten points less. We are not "
             "dressing that up. ", WHITE, True),
            ("In exchange: a third less volatility, a shallower drawdown, and "
             "zero policy breaches against thirty-seven. That is not "
             "outperformance — it is a different mandate, and for an "
             "institution with a mandated risk appetite that trade is the "
             "entire point.", MUTED, False),
        ],
        size=19, first=True,
    )

    half, gap = Inches(8.45), Inches(0.50)
    left = _box(slide, LEFT, Inches(7.10), half, Inches(3.4))
    _label(left, "Feasibility — it runs today", first=True)
    for text in (
        "Deployed and publicly reachable. Runs with no network and no API "
        "key: a market snapshot is committed to the repo.",
        "Six failure drills pass — no network, no key, deleted database, a "
        "−40% shock, approving a rejected candidate, weakening a threshold "
        "without a reason.",
        "Every threshold lives in config/policy.yaml. Changing one is "
        "versioned, requires a written reason, and is audited.",
    ):
        _rich(left, [("▸  ", GREEN, True), (text, MUTED, False)],
              size=18, space_after=7)

    right = _box(slide, LEFT + half + gap, Inches(7.10), half, Inches(3.4))
    _label(right, "Impact — and what this is not", first=True)
    _rich(right, [
        ("For a risk manager: ", WHITE, True),
        ("the reason an allocation was refused, in numbers — observed value "
         "against threshold — not a colour or a generic 'constraints "
         "violated'.", MUTED, False),
    ], size=18, space_after=7)
    _rich(right, [
        ("For an auditor: ", WHITE, True),
        ("what the system did, what it refused to do, and where a person "
         "stepped in — replayable from stored records.", MUTED, False),
    ], size=18, space_after=7)
    _rich(right, [
        ("Not: ", RED, True),
        ("a trading bot, connected to any brokerage, or a compliance "
         "product. Approval triggers a simulated rebalance. Thresholds are "
         "configurable demonstration values, not Basel/SEBI/RBI limits.",
         MUTED, False),
    ], size=18)


# ---------------------------------------------------------------------------
# SLIDE 7 — references
# ---------------------------------------------------------------------------


def _references(slide) -> None:
    tf = _box(slide, LEFT, TOP, WIDTH, BOTTOM - TOP)

    _label(tf, "Project links", first=True)
    for label, url in (
        ("Live dashboard", "https://capitalcontrol-engine.streamlit.app"),
        ("Project page", "https://pleaeswoekr.vercel.app"),
        ("Source code", "https://github.com/Adityarane012/CCE"),
    ):
        _rich(tf, [(f"{label}   ", WHITE, True), (url, GREEN, False)],
              size=23, space_after=9)

    _label(tf, "Documentation — the full 18-document specification is public")
    for label, url in (
        ("Safety invariants — the twelve claims, each with a test",
         "docs/10-RULES.md"),
        ("Architecture and the layer-dependency rules", "docs/02-ARCHITECTURE.md"),
        ("Every formula, with units and annualisation",
         "docs/08-FINANCIAL-METHODS.md"),
        ("Thresholds, control codes and breaker triggers",
         "docs/07-RISK-POLICY.md"),
        ("Build log — what actually went wrong, not just what was intended",
         "docs/IMPLEMENTATION-PLAN.md"),
    ):
        _rich(tf, [(f"{label}   ", MUTED, False), (url, DIM, False)],
              size=18, space_after=6)

    _label(tf, "Data and libraries")
    _p(
        tf,
        "jugaad-data (NSE / RBI market data) · CVXPY · NumPy · pandas · SciPy "
        "· Streamlit · Plotly · SQLite",
        size=18, colour=MUTED,
    )
    _p(
        tf,
        "Market data: NSE index and equity series, Sep 2023 – Aug 2026, "
        "committed as a cached snapshot for reproducibility.",
        size=16, colour=DIM, space_before=4,
    )

    note = _box(slide, LEFT, Inches(9.95), WIDTH, Inches(0.8))
    _p(
        note,
        "Decision-support prototype. Simulated execution only — no broker "
        "connection and no real orders. No guarantee of returns is made or "
        "implied.",
        size=16, colour=DIM, first=True,
    )


if __name__ == "__main__":
    raise SystemExit(main())
